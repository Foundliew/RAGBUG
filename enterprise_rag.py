import ray
import os
import faiss
import json
import numpy as np
import pandas as pd
import logging
import sys
import time
import re
import hashlib
import importlib.metadata
from sentence_transformers import SentenceTransformer
from unstructured.partition.auto import partition
from retrying import retry
from pathlib import Path
import psutil
import langdetect
import charset_normalizer
from logging.handlers import RotatingFileHandler
import socket
from fasttext import load_model
from queue import PriorityQueue
from dataclasses import dataclass
from typing import List
try:
    import magic_pdf
    from magic_pdf.data.dataset import PymuDocDataset
    from magic_pdf.model.doc_analyze_by_custom_model import doc_analyze
    from magic_pdf.config.enums import SupportedPdfParseMethod
    from magic_pdf.data.data_reader_writer import FileBasedDataWriter
except ImportError:
    magic_pdf = None
    logging.warning("magic_pdf module not found, falling back to unstructured for PDF processing")
from bs4 import BeautifulSoup
from pptx import Presentation
try:
    import markdown
except ImportError:
    markdown = None
    logging.warning("markdown module not found, falling back to plain text processing")
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None
    logging.warning("PyPDF2 module not found, PDF integrity checks will be skipped")

# Script version
SCRIPT_VERSION = "2025-05-20-optimized-v31"

# Configure logging with rotation
log_formatter = logging.Formatter("%(asctime)s [%(levelname)s] %(message)s")
vector_log_handler = RotatingFileHandler(
    "/app/vector_store.log", mode="a", maxBytes=10*1024*1024, backupCount=2
)
vector_log_handler.setFormatter(log_formatter)
debug_log_handler = RotatingFileHandler(
    "/app/debug.log", mode="a", maxBytes=10*1024*1024, backupCount=2
)
debug_log_handler.setFormatter(log_formatter)
stdout_handler = logging.StreamHandler(sys.stdout)
stdout_handler.setFormatter(log_formatter)

logging.basicConfig(
    level=logging.DEBUG,
    handlers=[vector_log_handler, debug_log_handler, stdout_handler]
)
logger = logging.getLogger(__name__)

# Environment
os.environ["TRANSFORMERS_OFFLINE"] = "1"
os.environ["HF_HUB_OFFLINE"] = "1"
os.environ["RAY_RUNTIME_ENV_TEMPORARY_REFERENCE_EXPIRATION_S"] = "3600"

# Load keyword replacements
KEYWORD_REPLACEMENTS_PATH = "/app/keyword_replacements.json"
def load_keyword_replacements():
    try:
        if os.path.exists(KEYWORD_REPLACEMENTS_PATH):
            with open(KEYWORD_REPLACEMENTS_PATH, "r", encoding="utf-8") as f:
                return json.load(f)
        else:
            logger.warning(f"Keyword replacements file {KEYWORD_REPLACEMENTS_PATH} not found, using default")
            return {"于生": "余生"}
    except Exception as e:
        logger.error(f"Failed to load keyword replacements: {e}", exc_info=True)
        return {"于生": "余生"}

KEYWORD_REPLACEMENTS = load_keyword_replacements()

def chunk_text(text, chunk_size=1000, overlap=200):
    """Split text into overlapping chunks with paragraph/sentence boundaries."""
    logger.debug(f"Chunking text, length: {len(text)}, chunk_size: {chunk_size}, overlap: {overlap}")
    if not text or len(text) <= chunk_size:
        return [text] if text else []
    
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        if end < len(text):
            paragraph_break = text.rfind('\n\n', start, end)
            if paragraph_break > start + chunk_size // 2:
                end = paragraph_break + 2
            else:
                sentence_break = max(
                    text.rfind('. ', start, end),
                    text.rfind('! ', start, end),
                    text.rfind('? ', start, end)
                )
                if sentence_break > start + chunk_size // 2:
                    end = sentence_break + 2
        
        chunk = text[start:end]
        if chunk.strip():
            chunks.append(chunk)
        start = end - overlap
    
    logger.info(f"Created {len(chunks)} chunks, sample: {chunks[0][:50] if chunks else 'None'}")
    return chunks

def test_log_writability():
    """Test if log files are writable."""
    logger.info("Testing log file writability")
    log_files = ["/app/vector_store.log", "/app/debug.log"]
    for log_file in log_files:
        try:
            if os.path.exists(log_file):
                if os.access(log_file, os.W_OK):
                    logger.info(f"Write permission confirmed for {log_file}")
                else:
                    logger.error(f"No write permission for {log_file}")
            else:
                logger.info(f"Log file {log_file} does not exist, attempting to create")
                with open(log_file, "a") as f:
                    f.write("")
                logger.info(f"Created {log_file}")
        except Exception as e:
            logger.error(f"Failed to test/create {log_file}: {e}", exc_info=True)
        finally:
            for handler in logger.handlers:
                handler.flush()

def test_network_connectivity(host, port, retries=5, delay=2):
    """Test network connectivity to Ray cluster with retries."""
    logger.info(f"Testing network connectivity to {host}:{port}")
    for attempt in range(retries):
        try:
            with socket.create_connection((host, port), timeout=5) as sock:
                logger.info(f"Successfully connected to {host}:{port} on attempt {attempt + 1}")
                return True
        except Exception as e:
            logger.warning(f"Failed to connect to {host}:{port} on attempt {attempt + 1}: {e}")
            time.sleep(delay)
    logger.error(f"Failed to connect to {host}:{port} after {retries} attempts")
    return False

def check_dependencies():
    """Verify critical dependencies and log versions with minimum version warnings."""
    logger.info(f"Checking dependencies for enterprise_rag_optimized.py version {SCRIPT_VERSION}")
    required = [
        ("torch", "2.0.0"),
        ("transformers", "4.30.0"),
        ("sentence-transformers", "2.2.0"),
        ("faiss-cpu", "1.7.0"),
        ("unstructured", "0.10.0"),
        ("pandas", "2.0.0"),
        ("psutil", "5.9.0"),
        ("langdetect", "1.0.9"),
        ("charset-normalizer", "3.0.0"),
        ("ray", "2.44.1"),
        ("retrying", "1.3.3"),
        ("beautifulsoup4", "4.12.0"),
        ("python-pptx", "0.6.21"),
    ]
    optional = [
        ("magic_pdf", "0.6.0"),
        ("markdown", "3.4.0"),
        ("PyPDF2", "3.0.0"),
    ]
    missing = []
    for module, min_version in required:
        try:
            version = importlib.metadata.version(module.replace("_", "-"))
            logger.info(f"Found {module} version {version}")
            if module == "faiss-cpu":
                import faiss
                logger.info("Successfully imported faiss")
            clean_version = re.match(r"^\d+\.\d+\.\d+", version).group(0) if re.match(r"^\d+\.\d+\.\d+", version) else version
            from packaging import version as version_parse
            if version_parse.parse(clean_version) < version_parse.parse(min_version):
                logger.warning(f"{module} version {version} is below recommended {min_version}")
        except (importlib.metadata.PackageNotFoundError, ImportError) as e:
            missing.append(f"{module}: {str(e)}")
            logger.error(f"Dependency missing: {module}: {str(e)}")
    for module, min_version in optional:
        try:
            version = importlib.metadata.version(module.replace("_", "-"))
            logger.info(f"Found optional {module} version {version}")
            clean_version = re.match(r"^\d+\.\d+\.\d+", version).group(0) if re.match(r"^\d+\.\d+\.\d+", version) else version
            from packaging import version as version_parse
            if version_parse.parse(clean_version) < version_parse.parse(min_version):
                logger.warning(f"{module} version {version} is below recommended {min_version}")
        except (importlib.metadata.PackageNotFoundError, ImportError) as e:
            logger.warning(f"Optional dependency {module} not found: {str(e)}")
    if missing:
        logger.error(f"Critical dependencies missing: {', '.join(missing)}. Please install them on all Ray nodes (e.g., pip install {' '.join(m.split(':')[0] for m in missing)})")
        raise ImportError(f"Missing dependencies: {', '.join(missing)}")
    logger.info("Dependency check completed")

def log_system_resources():
    """Log available CPU and memory resources."""
    cpu_count = psutil.cpu_count()
    cpu_percent = psutil.cpu_percent(interval=1)
    mem = psutil.virtual_memory()
    logger.info(
        f"System resources: CPUs={cpu_count}, CPU usage={cpu_percent}%, "
        f"Memory total={mem.total/1024/1024:.2f}MB, available={mem.available/1024/1024:.2f}MB"
    )

@dataclass
class Task:
    file_path: str
    chunk_count: int
    file_type: str
    priority: int  # Lower number = higher priority (e.g., small files)

@dataclass
class VectorStoreInfo:
    actor: 'VectorStore'
    current_load: int  # Number of tasks
    max_capacity: int  # Max simultaneous tasks
    processing_time_avg: float  # Average processing time per chunk (seconds)

class TaskScheduler:
    def __init__(self, vector_stores: List[VectorStoreInfo]):
        self.vector_stores = vector_stores
        self.task_queue = PriorityQueue()

    def add_task(self, task: Task):
        logger.debug(f"Adding task to queue: {task.file_path}, priority: {task.priority}")
        self.task_queue.put((task.priority, task))

    def assign_task(self):
        if self.task_queue.empty():
            logger.debug("Task queue is empty")
            return None
        # Select actor with lowest relative load
        min_load_store = min(
            self.vector_stores,
            key=lambda x: x.current_load / x.max_capacity,
            default=None
        )
        if min_load_store and min_load_store.current_load < min_load_store.max_capacity:
            priority, task = self.task_queue.get()
            min_load_store.current_load += 1
            logger.info(f"Assigned task {task.file_path} to actor_{min_load_store.actor.actor_id}")
            return min_load_store.actor, task
        logger.debug("No available actors with capacity")
        return None

    def complete_task(self, actor_id: int):
        for store in self.vector_stores:
            if store.actor.actor_id == actor_id:
                store.current_load -= 1
                logger.debug(f"Task completed on actor_{actor_id}, new load: {store.current_load}")
                break

@ray.remote(num_cpus=16, memory=32 * 1024 * 1024 * 1024)
class VectorStore:
    def __init__(self, cache_dir: str = "/mnt/nfs_doc_cache", actor_id: int = 0):
        logger.info(f"Starting VectorStore initialization, version {SCRIPT_VERSION}, actor_id: {actor_id}")
        self.cache_dir = cache_dir
        self.model_cache_dir = os.path.join("/models")
        self.actor_id = actor_id
        self.current_tasks = {}  # {file_path: start_time}
        self.processed_file_hashes = set()  # For deduplication
        try:
            os.makedirs(self.model_cache_dir, exist_ok=True)
            logger.info(f"Using cache directory: {self.cache_dir}")
        except Exception as e:
            logger.error(f"Failed to create cache directory {self.model_cache_dir}: {e}", exc_info=True)
            raise

        process = psutil.Process()
        mem_before = process.memory_info().rss / 1024 / 1024
        logger.info(f"Memory usage before init: {mem_before:.2f} MB")

        try:
            self.texts = []
            self.chunk_metadata = []
            self.file_data = {}
            self.dimension = 1024

            self.index = faiss.IndexHNSWFlat(self.dimension, 32)
            self.index.hnsw.efConstruction = 200
            self.index.hnsw.efSearch = 100
            logger.info("FAISS index initialized")

            # Disable NNPACK for old CPUs
            os.environ["USE_NNPACK"] = "0"

            @retry(stop_max_attempt_number=5, wait_fixed=5000)
            def load_model(model_name, model_class, local_path=None, **kwargs):
                logger.info(f"Attempting to load model: {model_name}, local_path: {local_path}")
                if local_path and os.path.exists(local_path):
                    logger.info(f"Loading model from local path: {local_path}")
                    return model_class(local_path, **kwargs)
                logger.info(f"Loading model from name: {model_name}")
                return model_class(model_name, **kwargs)

            logger.debug("Loading Chinese model")
            chinese_model_path = os.path.join(self.model_cache_dir, "BAAI/bge-large-zh-v1.5")
            logger.info(f"Chinese model path exists: {os.path.exists(chinese_model_path)}")
            try:
                self.chinese_model = load_model(
                    "BAAI/bge-large-zh-v1.5",
                    SentenceTransformer,
                    local_path=chinese_model_path,
                    cache_folder=self.model_cache_dir
                )
                logger.debug("Testing Chinese model")
                test_embedding = self.chinese_model.encode(["测试句子"])
                logger.info(f"Chinese model test successful, embedding shape: {test_embedding.shape}")
            except Exception as e:
                logger.error(f"Failed to load Chinese model: {e}", exc_info=True)
                raise

            logger.debug("Loading English model")
            english_model_path = os.path.join(self.model_cache_dir, "BAAI/bge-m3")
            logger.info(f"English model path exists: {os.path.exists(english_model_path)}")
            try:
                self.english_model = load_model(
                    "BAAI/bge-m3",
                    SentenceTransformer,
                    local_path=english_model_path,
                    cache_folder=self.model_cache_dir
                )
                logger.debug("Testing English model")
                test_embedding = self.english_model.encode(["Test sentence"])
                logger.info(f"English model test successful, embedding shape: {test_embedding.shape}")
            except Exception as e:
                logger.error(f"Failed to load English model: {e}", exc_info=True)
                raise

            self.index_path = os.path.join(self.cache_dir, f"vector_store_actor_{actor_id}.index")
            self.texts_path = os.path.join(self.cache_dir, f"texts_actor_{actor_id}.txt")
            self.metadata_path = os.path.join(self.cache_dir, f"metadata_actor_{actor_id}.txt")
            logger.info(f"Storage paths: index={self.index_path}, texts={self.texts_path}, metadata={self.metadata_path}")

            self.load()

            mem_after = process.memory_info().rss / 1024 / 1024
            logger.info(f"Memory usage after init: {mem_after:.2f} MB")

        except Exception as e:
            logger.error(f"VectorStore initialization failed: {e}", exc_info=True)
            raise
        finally:
            for handler in logger.handlers:
                handler.flush()

        logger.info("VectorStore initialization complete")

    def health_check(self):
        """Extended health check with vector count and task status"""
        try:
            vector_count = self.index.ntotal if hasattr(self.index, 'ntotal') else 0
            chunk_count = len(self.texts)
            file_count = len(set(meta.get('file_name', '') for meta in self.chunk_metadata))
            
            stuck_tasks = []
            current_time = time.time()
            for file_path, start_time in self.current_tasks.items():
                if current_time - start_time > 600:  # Increased from 300 to 600 seconds
                    stuck_tasks.append(f"{file_path} ({int(current_time - start_time)}s)")
            
            model_status = self.model_health_check()
            model_healthy = model_status.get("chinese_model") and model_status.get("english_model")
            
            health_info = {
                "status": "healthy" if model_healthy and not stuck_tasks else "unhealthy",
                "vector_count": vector_count,
                "chunk_count": chunk_count,
                "file_count": file_count,
                "base_index_size": vector_count,
                "memory_mb": psutil.Process().memory_info().rss / 1024 / 1024,
                "timestamp": time.time(),
                "actor_id": self.actor_id,
                "current_tasks": len(self.current_tasks),
                "stuck_tasks": stuck_tasks
            }
            
            if not model_healthy:
                health_info["error"] = "One or more embedding models are unhealthy"
            if stuck_tasks:
                health_info["error"] = f"Stuck tasks detected: {', '.join(stuck_tasks)}"
            
            logger.info(f"Health check: {json.dumps(health_info)}")
            return health_info
            
        except Exception as e:
            logger.error(f"Health check failed: {e}", exc_info=True)
            return {
                "status": "unhealthy",
                "error": str(e),
                "actor_id": self.actor_id,
                "base_index_size": 0
            }

    def model_health_check(self):
        """Check the operational status of the embedding models."""
        logger.info("Performing model health check")
        status = {"chinese_model": False, "english_model": False, "details": {}}
        try:
            start_time = time.time()
            test_embedding = self.chinese_model.encode(["测试句子"])
            elapsed_time = time.time() - start_time
            if test_embedding.shape == (1, self.dimension):
                status["chinese_model"] = True
                status["details"]["chinese_model"] = {
                    "embedding_shape": test_embedding.shape,
                    "sample_values": test_embedding[0][:5].tolist(),
                    "time_taken": elapsed_time
                }
                logger.info(f"Chinese model (BAAI/bge-large-zh-v1.5) healthy, shape: {test_embedding.shape}, time: {elapsed_time:.2f} seconds")
            else:
                logger.error(f"Chinese model failed: Invalid embedding shape {test_embedding.shape}")
        except Exception as e:
            logger.error(f"Chinese model health check failed: {e}", exc_info=True)
            status["details"]["chinese_model"] = {"error": str(e)}

        try:
            start_time = time.time()
            test_embedding = self.english_model.encode(["Test sentence"])
            elapsed_time = time.time() - start_time
            if test_embedding.shape == (1, self.dimension):
                status["english_model"] = True
                status["details"]["english_model"] = {
                    "embedding_shape": test_embedding.shape,
                    "sample_values": test_embedding[0][:5].tolist(),
                    "time_taken": elapsed_time
                }
                logger.info(f"English model (BAAI/bge-m3) healthy, shape: {test_embedding.shape}, time: {elapsed_time:.2f} seconds")
            else:
                logger.error(f"English model failed: Invalid embedding shape {test_embedding.shape}")
        except Exception as e:
            logger.error(f"English model health check failed: {e}", exc_info=True)
            status["details"]["english_model"] = {"error": str(e)}

        logger.info(f"Model health check result: {status}")
        return status

    @retry(stop_max_attempt_number=3, wait_fixed=2000)
    def load(self):
        """Load vector store from disk"""
        try:
            if os.path.exists(self.index_path) and os.path.exists(self.texts_path) and os.path.exists(self.metadata_path):
                logger.info(f"Loading existing vector store from {self.index_path}")
                
                self.index = faiss.read_index(self.index_path)
                
                with open(self.texts_path, 'r', encoding='utf-8') as f:
                    self.texts = [line.strip() for line in f if line.strip()]
                
                with open(self.metadata_path, 'r', encoding='utf-8') as f:
                    self.chunk_metadata = [json.loads(line) for line in f if line.strip()]
                
                logger.info(f"Loaded {len(self.texts)} documents with {self.index.ntotal} vectors")
            else:
                logger.info("No existing vector store found, creating new one")
                self.texts = []
                self.chunk_metadata = []
                self.index = faiss.IndexHNSWFlat(self.dimension, 32)
                self.index.hnsw.efConstruction = 200
                self.index.hnsw.efSearch = 100
            if len(self.texts) != len(self.chunk_metadata):
                logger.warning(f"Mismatch in loaded data: {len(self.texts)} texts vs {len(self.chunk_metadata)} metadata")
        except Exception as e:
            logger.error(f"Failed to load vector store: {e}", exc_info=True)
            self.texts = []
            self.chunk_metadata = []
            self.index = faiss.IndexHNSWFlat(self.dimension, 32)
            self.index.hnsw.efConstruction = 200
            self.index.hnsw.efSearch = 100

    @retry(stop_max_attempt_number=3, wait_fixed=2000)
    def save(self):
        """Save vector store to disk"""
        logger.info(f"Saving vector store to {self.index_path}")
        try:
            faiss.write_index(self.index, self.index_path)
            
            with open(self.texts_path, "w", encoding="utf-8") as f:
                f.write("\n".join(self.texts) + "\n")
            with open(self.metadata_path, "w", encoding="utf-8") as f:
                for meta in self.chunk_metadata:
                    f.write(json.dumps(meta, ensure_ascii=False) + "\n")
                summary = {
                    "total_vectors": self.index.ntotal,
                    "files": {k: v.get("vector_count", 0) for k, v in self.file_data.items()}
                }
                f.write(json.dumps(summary, ensure_ascii=False) + "\n")
            if not os.path.exists(self.index_path) or not os.path.exists(self.texts_path) or not os.path.exists(self.metadata_path):
                logger.error("One or more save files were not created")
                raise IOError("Save operation failed")
            logger.info(f"Saved {len(self.texts)} documents with {self.index.ntotal} vectors")
        except Exception as e:
            logger.error(f"Failed to save data: {e}", exc_info=True)
            raise

    def detect_encoding(self, file_path):
        logger.debug(f"Detecting encoding for {file_path}")
        try:
            with open(file_path, "rb") as f:
                content = f.read()
            result = charset_normalizer.detect(content)
            encoding = result["encoding"] or "utf-8"
            logger.info(f"Detected encoding for {file_path}: {encoding}")
            return encoding
        except Exception as e:
            logger.error(f"Encoding detection failed for {file_path}: {e}", exc_info=True)
            return "utf-8"

    def apply_keyword_replacements(self, content):
        """Apply keyword replacements from configuration."""
        logger.debug(f"Applying keyword replacements to content, length: {len(content)}")
        for old, new in KEYWORD_REPLACEMENTS.items():
            content = content.replace(old, new)
        logger.debug(f"Content after replacements, length: {len(content)}")
        return content

    def clean_text(self, text):
        logger.debug(f"Cleaning text, original length: {len(text)}, sample: {text[:50] if text else 'None'}")
        text = re.sub(r"\s+", " ", text).strip()
        text = re.sub(r"[\x00-\x08\x0B\x0C\x0E-\x1F]", "", text)
        text = self.apply_keyword_replacements(text)
        logger.debug(f"Cleaned text length: {len(text)}, sample: {text[:50] if text else 'None'}")
        return text

    def split_text(self, text, chunk_size=500, chunk_overlap=0.1):
        """Split text based on Markdown structure or plain text."""
        logger.debug(f"Splitting text, length: {len(text)}, sample: {text[:50] if text else 'None'}")
        if not text or len(text.strip()) < 5:
            logger.warning("Text too short or empty, skipping chunking")
            return []

        overlap = int(chunk_size * chunk_overlap)
        if markdown:
            sections = re.split(r'(^#{1,6}\s.*$|\n+)', text, flags=re.MULTILINE)
            sections = [s.strip() for s in sections if s.strip()]
            chunks = []
            current_chunk = ""
            for section in sections:
                section_length = len(section)
                if len(current_chunk) + section_length <= chunk_size:
                    current_chunk += section + "\n\n"
                else:
                    if current_chunk:
                        chunks.append(current_chunk.strip())
                    current_chunk = section + "\n\n"
                    if chunks and overlap > 0:
                        last_chunk = chunks[-1][-overlap:] if len(chunks[-1]) > overlap else chunks[-1]
                        current_chunk = last_chunk + current_chunk
            if current_chunk:
                chunks.append(current_chunk.strip())
            if not chunks and text.strip():
                chunks.append(text)
        else:
            chunks = chunk_text(text, chunk_size=chunk_size, overlap=overlap)

        for i, chunk in enumerate(chunks):
            logger.debug(f"Chunk {i+1}/{len(chunks)}, length: {len(chunk)}, sample: {chunk[:50] if chunk else 'None'}")
        logger.info(f"Split text into {len(chunks)} chunks, sample: {chunks[0][:50] if chunks else 'None'}")
        return chunks

    def get_chunk_size(self, file_size, content_type="text"):
        """Determine chunk size based on file size and content type."""
        logger.debug(f"Determining chunk size for file size: {file_size}, content_type: {content_type}")
        if content_type == "markdown" and markdown:
            return 1000
        if file_size < 10 * 1024:
            return 500
        elif file_size > 1024 * 1024:
            return 1500
        return 1000

    def process_pdf(self, file_path, file_content=None):
        """Process PDF file and extract text correctly."""
        filename = os.path.basename(file_path)
        logger.info(f"PDF processing strategy: {'magic_pdf' if magic_pdf else 'unstructured'}")
        content = ""
        metadata = {
            "file_path": file_path,
            "file_type": "pdf",
            "tables": [],
            "images": [],
            "sections": [],
            "type": "text" if not markdown else "markdown"
        }
        
        if magic_pdf:
            try:
                image_dir = os.path.join(self.cache_dir, "images")
                os.makedirs(image_dir, exist_ok=True)
                image_writer = FileBasedDataWriter(image_dir)
                logger.info(f"Initialized FileBasedDataWriter for {image_dir}")
                
                if file_content and isinstance(file_content, bytes):
                    pdf_bytes = file_content
                else:
                    with open(file_path, "rb") as f:
                        pdf_bytes = f.read()
                ds = PymuDocDataset(pdf_bytes)
                is_ocr = ds.classify() == SupportedPdfParseMethod.OCR
                logger.info(f"PDF requires OCR: {is_ocr}")
                
                if is_ocr:
                    infer_result = ds.apply(doc_analyze, ocr=True)
                    content = infer_result.pipe_ocr_mode(image_writer).get_markdown(os.path.basename(image_dir))
                else:
                    infer_result = ds.apply(doc_analyze, ocr=False)
                    content = infer_result.pipe_txt_mode(image_writer).get_markdown(os.path.basename(image_dir))
                
                content = self.clean_text(content)
                logger.info(f"Extracted Markdown content from MinerU, length: {len(content)}")
            except Exception as e:
                logger.warning(f"Magic PDF processing failed: {e}, falling back to unstructured")
        
        if not content or content.startswith("[Error:"):
            strategies = ["hi_res", "ocr_only", "fast"]
            for strategy in strategies:
                try:
                    elements = partition(filename=file_path, strategy=strategy, include_page_breaks=True)
                    content_parts = []
                    image_dir = os.path.join(self.cache_dir, "images")
                    os.makedirs(image_dir, exist_ok=True)
                    
                    for el in elements:
                        text = str(el)
                        if hasattr(el, "category") and el.category == "Table":
                            table_text = el.metadata.get("text_as_html", text)
                            try:
                                df = pd.read_html(table_text)[0]
                                if markdown:
                                    header_row = "| " + " | ".join(map(str, df.columns)) + " |"
                                    separator = "| " + " | ".join(["---" for _ in df.columns]) + " |"
                                    rows = ["| " + " | ".join(map(str, row)) + " |" for _, row in df.iterrows()]
                                    table_md = "\n".join([header_row, separator] + rows)
                                else:
                                    table_md = df.to_string(index=False)
                                metadata["tables"].append(table_md)
                                content_parts.append(table_md)
                            except:
                                content_parts.append(f"Table: {text}")
                        elif hasattr(el, "category") and el.category == "Image":
                            img_desc = el.metadata.get("text", "No description available")
                            image_path = os.path.join(image_dir, f"{hashlib.md5(img_desc.encode()).hexdigest()}_{filename}")
                            metadata["images"].append({"path": image_path, "description": img_desc})
                            if markdown:
                                content_parts.append(f"![Image]({image_path})\nDescription: {img_desc}")
                            else:
                                content_parts.append(f"Image: {image_path}\nDescription: {img_desc}")
                        else:
                            content_parts.append(text)
                    
                    content = "\n\n".join(content_parts)
                    content = self.clean_text(content)
                    logger.info(f"Processed PDF with strategy {strategy}, content length: {len(content)}")
                    break
                except Exception as e:
                    logger.warning(f"Failed to process PDF with strategy {strategy}: {e}")
            else:
                logger.error(f"All strategies failed for {file_path}")
                return f"[Error: Failed to process {filename}]", metadata
        
        if markdown:
            section_matches = re.finditer(r'^#{1,6}\s(.*)$', content, re.MULTILINE)
            for match in section_matches:
                metadata["sections"].append({
                    "title": match.group(1),
                    "start_pos": match.start()
                })
        
        logger.info(f"Detected {len(metadata['tables'])} tables and {len(metadata['images'])} images in {filename}")
        return content, metadata

    def process_pdf_batch(self, batch_bytes, file_path, metadata):
        """Process a batch of PDF pages."""
        filename = os.path.basename(file_path)
        try:
            start_time = time.perf_counter()
            content = ""
            if magic_pdf:
                try:
                    image_dir = os.path.join(self.cache_dir, "images")
                    os.makedirs(image_dir, exist_ok=True)
                    image_writer = ["FileBasedDataWriter(image_dir)"]
                    ds = PymuDocDataset(batch_bytes)
                    is_ocr = ds.classify() == SupportedPdfParseMethod.OCR
                    if is_ocr:
                        infer_result = ds.apply(doc_analyze, ocr=True)
                        content = infer_result.pipe_ocr_mode(image_writer).get_markdown(os.path.basename(image_dir))
                    else:
                        infer_result = ds.apply(doc_analyze, ocr=False)
                        content = infer_result.pipe_txt_mode(image_writer).get_markdown(os.path.basename(image_dir))
                    content = self.clean_text(content)
                    logger.debug(f"Processed batch for {filename}, length: {len(content)}, time: {time.perf_counter() - start_time:.2f}s")
                except Exception as e:
                    logger.warning(f"magic_pdf batch failed for {filename}: {e}")
            if not content or content.startswith("[Error:"):
                parts = []
                elements = partition(file_content=batch_bytes, strategy="ocr_only", include_page_breaks=False)
                for el in elements:
                    text = str(el)
                    if hasattr(el, "category") and el.category == "Table":
                        parts.append(f"Table: {text}")
                    elif hasattr(el, "text"):
                        parts.append(text)
                content = "\n\n".join(parts)
                content = self.clean_text(content)
            debug_path = os.path.join(self.cache_dir, f"debug_{filename}.txt")
            with open(debug_path, "a", encoding="utf-8") as f:
                f.write(f"Batch processed: {len(content)} chars\n{content}\n\n")
            return content
        except Exception as e:
            logger.error(f"Batch processing failed for {filename}: {e}")
            return ""

    def process_markdown(self, file_path, file_content=None):
        """Process Markdown file with proper header handling."""
        logger.info(f"Processing Markdown file: {file_path}")
        try:
            if file_content and isinstance(file_content, bytes):
                content = file_content.decode('utf-8', errors='replace')
            else:
                encoding = self.detect_encoding(file_path)
                try:
                    with open(file_path, "r", encoding=encoding) as f:
                        content = f.read()
                except Exception as e:
                    logger.warning(f"Failed to read {file_path} with encoding {encoding}: {e}")
                    with open(file_path, "rb") as f:
                        content = f.read().decode("utf-8", errors="ignore")
            
            content = re.sub(r'(#+\s+[^\n]+)([^#\n])', r'\1\n\2', content)
            content = re.sub(r'(#+\s+[^\n]+\n)([^#\n])', r'\1\n\2', content)
            
            headers = []
            for match in re.finditer(r'^(#+)\s+(.+?)$', content, re.MULTILINE):
                level = len(match.group(1))
                text = match.group(2).strip()
                headers.append({"level": level, "text": text})
            
            if markdown:
                html = markdown.markdown(content)
                soup = BeautifulSoup(html, 'html.parser')
                extracted_text = soup.get_text('\n\n')
            else:
                extracted_text = content
            
            extracted_text = self.clean_text(extracted_text)
            metadata = {
                "file_path": file_path,
                "file_type": "markdown",
                "headers": headers,
                "tables": [],
                "images": [],
                "type": "markdown" if markdown else "text",
                "sections": [{"title": h["text"], "start_pos": content.find(h["text"])} for h in headers]
            }
            
            logger.info(f"Processed Markdown, content length: {len(extracted_text)}")
            return extracted_text, metadata
        
        except Exception as e:
            logger.error(f"Error processing Markdown {file_path}: {e}", exc_info=True)
            return f"[Error: Failed to process {os.path.basename(file_path)}]", {"error": str(e), "file_path": file_path, "file_type": "markdown"}

    def process_file(self, file_path):
        filename = os.path.basename(file_path)
        file_size = os.path.getsize(file_path)
        content_type = "markdown" if markdown and filename.lower().endswith(".md") else "text"
        chunk_size = self.get_chunk_size(file_size, content_type=content_type)
        logger.info(f"Processing file: {file_path} (size: {file_size} bytes, chunk_size: {chunk_size}, content_type: {content_type})")

        supported_extensions = {
            ".txt", ".md", ".csv", ".xlsx", ".xls", ".json", ".pdf", ".docx",
            ".jpg", ".jpeg", ".png", ".ppt", ".pptx", ".html"
        }

        if not any(filename.lower().endswith(ext) for ext in supported_extensions):
            logger.warning(f"Unsupported file type: {filename}")
            return "", None

        content = ""
        metadata = {"tables": [], "header": None, "type": content_type, "sections": [], "images": []}

        try:
            if filename.lower().endswith(".pdf") and PyPDF2:
                try:
                    with open(file_path, "rb") as f:
                        PyPDF2.PdfReader(f)
                    logger.info(f"PDF {filename} is valid")
                except Exception as e:
                    logger.error(f"Invalid PDF {filename}: {e}")
                    return f"[Error: Invalid PDF file {filename}]", metadata

            if filename.lower().endswith((".txt", ".md")):
                if filename.lower().endswith(".md"):
                    content, metadata = self.process_markdown(file_path)
                else:
                    encoding = self.detect_encoding(file_path)
                    try:
                        with open(file_path, "r", encoding=encoding) as f:
                            content = f.read()
                    except Exception as e:
                        logger.warning(f"Failed to read {file_path} with encoding {encoding}: {e}")
                        with open(file_path, "rb") as f:
                            content = f.read().decode("utf-8", errors="ignore")
                    content = self.clean_text(content)
                    metadata = {"file_path": file_path, "file_type": "text", "tables": [], "header": None, "type": "text", "sections": [], "images": []}
                logger.debug(f"Extracted content from {filename}, length: {len(content)}")

            elif filename.lower().endswith(".csv"):
                encoding = self.detect_encoding(file_path)
                delimiters = [",", ";", "\t"]
                for delimiter in delimiters:
                    try:
                        df = pd.read_csv(
                            file_path,
                            encoding=encoding,
                            delimiter=delimiter,
                            on_bad_lines="skip",
                            low_memory=False
                        )
                        logger.info(f"Successfully parsed CSV with delimiter: {delimiter}")
                        metadata["header"] = df.columns.tolist()
                        if markdown:
                            header_row = "| " + " | ".join(map(str, df.columns)) + " |"
                            separator = "| " + " | ".join(["---" for _ in df.columns]) + " |"
                            rows = ["| " + " | ".join(map(str, row)) + " |" for _, row in df.iterrows()]
                            content = "\n".join([header_row, separator] + rows)
                        else:
                            header_str = ", ".join(map(str, df.columns))
                            content = f"Header: {header_str}\n{df.to_string(index=False)}"
                        metadata["tables"].append(content)
                        content = self.clean_text(content)
                        logger.debug(f"Processed CSV, length: {len(content)}")
                        break
                    except Exception as e:
                        logger.warning(f"Failed to parse CSV with delimiter {delimiter}: {e}")
                else:
                    logger.error(f"All delimiter attempts failed for {file_path}")
                    return f"[Error: Failed to process {filename}]", metadata

            elif filename.lower().endswith((".xlsx", ".xls")):
                try:
                    df = pd.read_excel(file_path)
                    metadata["header"] = df.columns.tolist()
                    if markdown:
                        header_row = "| " + " | ".join(map(str, df.columns)) + " |"
                        separator = "| " + " | ".join(["---" for _ in df.columns]) + " |"
                        rows = ["| " + " | ".join(map(str, row)) + " |" for _, row in df.iterrows()]
                        content = "\n".join([header_row, separator] + rows)
                    else:
                        header_str = ", ".join(map(str, df.columns))
                        content = f"Header: {header_str}\n{df.to_string(index=False)}"
                    metadata["tables"].append(content)
                    content = self.clean_text(content)
                    logger.debug(f"Processed Excel, length: {len(content)}")
                except Exception as e:
                    logger.error(f"Failed to process Excel file {file_path}: {e}")
                    return f"[Error: Failed to process {filename}]", metadata

            elif filename.lower().endswith(".json"):
                try:
                    with open(file_path, "r", encoding=self.detect_encoding(file_path)) as f:
                        data = json.load(f)
                    if markdown:
                        content = markdown.markdown(json.dumps(data, ensure_ascii=False, indent=2))
                    else:
                        content = json.dumps(data, ensure_ascii=False, indent=2)
                    content = self.clean_text(content)
                    logger.debug(f"Processed JSON, length: {len(content)}")
                except Exception as e:
                    logger.error(f"Failed to process JSON file {file_path}: {e}")
                    return f"[Error: Failed to process {filename}]", metadata

            elif filename.lower().endswith((".jpg", ".jpeg", ".png")):
                try:
                    elements = partition(filename=file_path, strategy="hi_res")
                    content_parts = []
                    image_dir = os.path.join(self.cache_dir, "images")
                    os.makedirs(image_dir, exist_ok=True)
                    for el in elements:
                        if hasattr(el, "text") and el.text:
                            text = str(el)
                            if hasattr(el, "category") and el.category == "Image":
                                img_desc = el.metadata.get("text", "No description available")
                                image_path = os.path.join(image_dir, f"{hashlib.md5(img_desc.encode()).hexdigest()}_{filename}")
                                metadata["images"].append({"path": image_path, "description": img_desc})
                                logger.info(f"Saved image metadata for {filename} at {image_path}")
                                if markdown:
                                    content_parts.append(f"![Image Description]({image_path})\nDescription: {img_desc}")
                                else:
                                    content_parts.append(f"Image: {image_path}\nDescription: {img_desc}")
                            else:
                                content_parts.append(text)
                    content = "\n\n".join(content_parts) or f"Image: {filename}\nNo OCR text available"
                    content = self.clean_text(content)
                    logger.debug(f"Processed image, length: {len(content)}")
                except Exception as e:
                    logger.error(f"Failed to process image file {file_path}: {e}")
                    return f"[Error: Failed to process {filename}]", metadata

            elif filename.lower().endswith((".ppt", ".pptx")):
                try:
                    prs = Presentation(file_path)
                    content_parts = []
                    for slide_idx, slide in enumerate(prs.slides):
                        slide_content = f"Slide {slide_idx + 1}\n"
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                slide_content += shape.text + "\n"
                        content_parts.append(slide_content)
                    content = "\n\n".join(content_parts)
                    content = self.clean_text(content)
                    logger.debug(f"Processed PPT, length: {len(content)}")
                except Exception as e:
                    logger.error(f"Failed to process PPT file {file_path}: {e}")
                    return f"[Error: Failed to process {filename}]", metadata

            elif filename.lower().endswith(".html"):
                try:
                    with open(file_path, "r", encoding=self.detect_encoding(file_path)) as f:
                        soup = BeautifulSoup(f, "html.parser")
                    if markdown:
                        content = markdown.markdown(soup.get_text())
                    else:
                        content = soup.get_text()
                    tables = soup.find_all("table")
                    for table in tables:
                        try:
                            df = pd.read_html(str(table))[0]
                            if markdown:
                                header_row = "| " + " | ".join(map(str, df.columns)) + " |"
                                separator = "| " + " | ".join(["---" for _ in df.columns]) + " |"
                                rows = ["| " + " | ".join(map(str, row)) + " |" for _, row in df.iterrows()]
                                table_md = "\n".join([header_row, separator] + rows)
                            else:
                                table_md = df.to_string(index=False)
                            metadata["tables"].append(table_md)
                            content += f"\n\n{table_md}"
                        except Exception as e:
                            logger.warning(f"Failed to convert HTML table: {e}")
                    content = self.clean_text(content)
                    logger.debug(f"Processed HTML, length: {len(content)}")
                except Exception as e:
                    logger.error(f"Failed to process HTML file {file_path}: {e}")
                    return (f"[Error: Failed to processobservation failed: {filename}]", metadata)

            if content_type == "markdown" and markdown:
                section_matches = re.finditer(r'^#{1,6}\s(.*)$', content, re.MULTILINE)
                for match in section_matches:
                    metadata["sections"].append({
                        "title": match.group(1),
                        "start_pos": match.start()
                    })

            logger.info(f"Detected {len(metadata['tables'])} tables and {len(metadata['images'])} images in {filename}")
            return content, metadata
        except Exception as e:
            logger.error(f"Unexpected error processing {file_path}: {e}", exc_info=True)
            return f"[Error: Failed to process {filename}]", metadata

    def get_file_hash(self, file_path):
        """Calculate file hash for deduplication."""
        try:
            with open(file_path, "rb") as f:
                return hashlib.md5(f.read()).hexdigest()
        except Exception as e:
            logger.warning(f"Failed to hash {file_path}: {e}")
            return None

    @retry(stop_max_attempt_number=3, wait_fixed=2000)
    def add_document(self, file_path, file_content=None, file_type=None):
        """Add document to vector store with per-chunk vectorization"""
        filename = os.path.basename(file_path)
        logger.info(f"Adding document: {file_path} of type {file_type}")
        
        try:
            start_time = time.time()
            self.current_tasks[file_path] = start_time
            
            # Deduplication check
            file_hash = self.get_file_hash(file_path) if os.path.exists(file_path) else None
            if file_hash and file_hash in self.processed_file_hashes:
                logger.info(f"Skipped duplicate file: {filename} (hash: {file_hash})")
                del self.current_tasks[file_path]
                return 0

            if not file_type:
                file_type = os.path.splitext(filename)[1].lower().lstrip('.')
                if file_type == 'md':
                    file_type = 'markdown'
                elif file_type not in ['text', 'csv', 'xlsx', 'xls', 'json', 'pdf', 'docx', 'jpg', 'jpeg', 'png', 'ppt', 'pptx', 'html']:
                    file_type = 'text'

            if filename in self.file_data:
                logger.info(f"File {filename} already processed")
                del self.current_tasks[file_path]
                return 0

            if file_type == 'pdf':
                # Parallel batch processing for PDFs
                pdf_bytes = file_content if isinstance(file_content, bytes) else open(file_path, "rb").read()
                batch_size_bytes = 10 * 1024 * 1024  # 10MB per batch
                batches = [pdf_bytes[i:i+batch_size_bytes] for i in range(0, len(pdf_bytes), batch_size_bytes)]
                content_parts = []
                metadata = {
                    "file_path": file_path,
                    "file_type": "pdf",
                    "tables": [],
                    "images": [],
                    "sections": [],
                    "type": "text" if not markdown else "markdown"
                }
                @ray.remote(num_cpus=2)
                def process_batch(self, batch, file_path, metadata):
                    return self.process_pdf_batch(batch, file_path, metadata)

                futures = [process_batch.remote(self, batch, file_path, metadata) for batch in batches]
                content_parts = [ray.get(f) for f in futures]
                content = "\n\n".join([part for part in content_parts if part])
                if not content or content.startswith("[Error:"):
                    content, metadata = self.process_pdf(file_path, file_content)
            elif file_type == 'markdown':
                content, metadata = self.process_markdown(file_path, file_content)
            else:
                content, metadata = self.process_file(file_path)
            
            if not content or content.startswith("[Error:"):
                logger.warning(f"Skipping vectorization for {file_path} due to processing error: {content}")
                del self.current_tasks[file_path]
                return 0
            
            self.file_data[filename] = metadata or {
                "tables": [], "header": None, "type": file_type, "sections": [], "images": []
            }

            chunk_size = self.get_chunk_size(os.path.getsize(file_path) if os.path.exists(file_path) else len(content), content_type=file_type)
            logger.debug(f"Splitting text into chunks, chunk_size: {chunk_size}")
            chunk_start_time = time.time()
            text_chunks = self.split_text(content, chunk_size=chunk_size)
            chunk_elapsed_time = time.time() - chunk_start_time
            logger.info(f"Split text into {len(text_chunks)} chunks, time: {chunk_elapsed_time:.2f}s, sample: {text_chunks[0][:50] if text_chunks else 'None'}")
            
            if not text_chunks:
                logger.warning(f"No chunks generated for {file_path}")
                del self.current_tasks[file_path]
                return 0

            chunks = []
            for i, text in enumerate(text_chunks):
                chunk = {
                    "content": text,
                    "file_name": filename,
                    "metadata": metadata,
                    "section": f"chunk_{i+1}",
                    "upload_timestamp": time.strftime("%Y-%m-%d %H:%M:%S")
                }
                chunks.append(chunk)

            logger.debug(f"Encoding {len(chunks)} chunks for {filename}")
            encode_start_time = time.time()
            batch_size = 100
            embeddings = []
            for i in range(0, len(chunks), batch_size):
                batch = [chunk["content"] for chunk in chunks[i:i+batch_size]]
                try:
                    batch_start = time.perf_counter()
                    batch_embedding = self.encode_batch(batch)
                    embeddings.extend(batch_embedding)
                    logger.debug(f"Encoded batch {i//batch_size+1} for {filename}, count: {len(batch_embedding)}, time: {time.perf_counter()-batch_start:.2f}s")
                except Exception as e:
                    logger.error(f"Batch {i//batch_size+1} encoding failed for {filename}: {e}")
                    continue
            embeddings = np.vstack(embeddings) if embeddings else np.array([])
            encode_elapsed_time = time.time() - encode_start_time
            logger.info(f"Encoded {len(embeddings)} chunks, embedding shape: {embeddings.shape}, time: {encode_elapsed_time:.2f}s")

            if embeddings.size == 0 or embeddings.shape[0] != len(chunks):
                logger.error(f"Embedding generation failed or mismatch: expected {len(chunks)} embeddings, got {embeddings.shape[0] if embeddings.size else 0}")
                del self.current_tasks[file_path]
                return 0

            logger.debug(f"Adding {len(chunks)} embeddings to FAISS index for {filename}")
            index_start_time = time.time()
            self.index.add(embeddings)
            index_elapsed_time = time.time() - index_start_time
            logger.info(f"Added {len(chunks)} embeddings to FAISS index, time: {index_elapsed_time:.2f}s")

            start_idx = len(self.texts)
            for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
                text = chunk["content"]
                chunk_id = hashlib.md5(text.encode("utf-8")).hexdigest()
                chunk_metadata = {
                    "file_name": filename,
                    "chunk_id": chunk_id,
                    "chunk_index": start_idx + i,
                    "metadata": chunk.get("metadata", {
                        "tables": [], 
                        "header": None, 
                        "type": file_type,
                        "sections": [], 
                        "images": []
                    }),
                    "type": file_type,
                    "section": chunk.get("section", None),
                    "upload_timestamp": chunk.get("upload_timestamp", time.strftime("%Y-%m-%d %H:%M:%S"))
                }
                self.texts.append(text)
                self.chunk_metadata.append(chunk_metadata)
                logger.debug(f"Stored chunk {i+1}/{len(chunks)} for {filename}, chunk_id: {chunk_id}")

            self.file_data[filename].update({
                "filename": filename,
                "upload_timestamp": time.strftime("%Y-%m-%d %H:%M:%S"),
                "chunk_count": len(chunks),
                "vector_count": len(embeddings),
                "duration": time.time() - start_time
            })

            if file_hash:
                self.processed_file_hashes.add(file_hash)

            logger.debug(f"Saving index, texts, and metadata")
            save_start_time = time.time()
            self.save()
            save_elapsed_time = time.time() - save_start_time
            logger.info(f"Saved index, texts, and metadata, time: {save_elapsed_time:.2f}s")
            
            total_elapsed_time = time.time() - start_time
            logger.info(f"Successfully added {len(embeddings)} vectors for {filename}, total time: {total_elapsed_time:.2f}s")
            del self.current_tasks[file_path]
            return len(embeddings)
        
        except Exception as e:
            logger.error(f"Error adding document {file_path}: {e}", exc_info=True)
            del self.current_tasks[file_path]
            return 0

    @retry(stop_max_attempt_number=3, wait_fixed=2000)
    def add_file(self, file_path, content=None, metadata=None):
        """Compatibility method for older add_file"""
        return self.add_document(file_path, file_content=content, file_type=None)

    def process_existing_files(self, directory):
        supported_extensions = {".txt", ".md", ".csv", ".xlsx", ".xls", ".json",
                               ".pdf", ".docx", ".jpg", ".jpeg", ".png", ".ppt", ".pptx", ".html"}
        file_paths = [
            os.path.join(directory, f) for f in os.listdir(directory)
            if any(f.lower().endswith(ext) for ext in supported_extensions)
        ]
        logger.info(f"Found {len(file_paths)} files to process in {directory}")

        for file_path in file_paths:
            content, metadata = self.process_file(file_path)
            if content and not content.startswith("[Error:"):
                self.add_document(file_path, content, metadata.get("file_type"))
            else:
                logger.warning(f"Skipping {file_path} due to processing error")

        self.save()

    def classify_query_type(self, query):
        """Classify query type for dynamic adjustment."""
        query_lower = query.lower()
        if "summarize" in query_lower or "summary" in query_lower:
            return "summarization"
        elif "title" in query_lower or "generate title" in query_lower:
            return "title_generation"
        elif any(term in query_lower for term in ["technical", "wireless", "network", "hardware"]):
            return "technical"
        return "general"

    def encode_batch(self, texts, query_type="general"):
        try:
            logger.info(f"Encoding batch, text count: {len(texts)}, sample: {texts[0][:50] if texts else 'None'}")
            if not texts:
                logger.warning("Empty text batch provided")
                return np.array([])
            ft_model = load_model("/models/lid.176.bin")
            lang_counts = {}
            sample_size = min(5, len(texts))
            for text in texts[:sample_size]:
                if not text.strip():
                    continue
                prediction = ft_model.predict(text.replace("\n", " "), k=1)
                lang = prediction[0][0].replace("__label__", "")
                lang_counts[lang] = lang_counts.get(lang, 0) + 1
            dominant_lang = max(lang_counts, key=lang_counts.get, default="en")
            logger.info(f"Detected dominant language: {dominant_lang}, counts: {lang_counts}")

            if len(lang_counts) > 1 or dominant_lang not in ["en", "zh"]:
                logger.warning(f"Ambiguous language detection: {lang_counts}. Falling back to per-chunk detection.")
                chunk_langs = []
                for text in texts:
                    if text.strip():
                        prediction = ft_model.predict(text.replace("\n", " "), k=1)
                        chunk_lang = prediction[0][0].replace("__label__", "")
                        chunk_langs.append(chunk_lang)
                if chunk_langs:
                    dominant_lang = max(set(chunk_langs), key=chunk_langs.count, default="en")
                    logger.info(f"Per-chunk dominant language: {dominant_lang}")

            model = self.chinese_model if dominant_lang == "zh" else self.english_model
            model_name = "BAAI/bge-large-zh-v1.5" if dominant_lang == "zh" else "BAAI/bge-m3"
            logger.info(f"Selected model: {model_name} for batch_size: {len(texts)}")

            normalize_embeddings = query_type != "title_generation"
            max_batch_size = 500 if query_type != "technical" else 200

            embeddings = []
            start_time = time.time()
            for i in range(0, len(texts), max_batch_size):
                batch = texts[i:i + max_batch_size]
                logger.debug(f"Processing batch {i//max_batch_size + 1}, size: {len(batch)}")
                batch_start_time = time.time()
                batch_embeddings = model.encode(
                    batch,
                    batch_size=len(batch),
                    normalize_embeddings=normalize_embeddings,
                    show_progress_bar=False
                )
                batch_elapsed_time = time.time() - batch_start_time
                logger.info(f"Encoded batch {i//max_batch_size + 1}, time: {batch_elapsed_time:.2f} seconds")
                if batch_embeddings.shape[1] != self.dimension:
                    logger.error(f"Invalid embedding shape: {batch_embeddings.shape}")
                    raise ValueError(f"Invalid embedding shape: {batch_embeddings.shape}")
                embeddings.append(batch_embeddings)

            embeddings = np.vstack(embeddings) if embeddings else np.array([])
            total_elapsed_time = time.time() - start_time
            logger.info(f"Completed encoding {len(texts)} texts, shape: {embeddings.shape}, total time: {total_elapsed_time:.2f}s, sample: {embeddings[0][:5] if embeddings.size else 'None'}")
            return embeddings
        except Exception as e:
            logger.error(f"Batch encoding failed: {e}", exc_info=True)
            raise

    def add_text_chunks(self, chunks: list):
        """Add preprocessed text chunks to the vector store."""
        logger.info(f"Adding {len(chunks)} text chunks")
        start_time = time.time()
        try:
            texts = [chunk["content"] for chunk in chunks]
            if not texts:
                logger.warning("No text chunks provided")
                return 0

            logger.debug(f"Encoding {len(texts)} chunks")
            encode_start_time = time.time()
            embeddings = []
            for i, text in enumerate(texts):
                try:
                    embedding = self.encode_batch([text])
                    if embedding.shape[0] != 1:
                        logger.error(f"Invalid embedding for chunk {i+1}: shape {embedding.shape}")
                        continue
                    embeddings.append(embedding[0])
                except Exception as e:
                    logger.error(f"Failed to encode chunk {i+1}: {e}")
                    continue
            embeddings = np.vstack(embeddings) if embeddings else np.array([])
            encode_elapsed_time = time.time() - encode_start_time
            logger.info(f"Encoded {len(embeddings)} chunks, embedding shape: {embeddings.shape}, time: {encode_elapsed_time:.2f}s")

            if embeddings.size == 0:
                logger.warning("No embeddings generated for text chunks")
                return 0

            logger.debug(f"Adding {len(embeddings)} embeddings to FAISS index")
            index_start_time = time.time()
            self.index.add(embeddings)
            index_elapsed_time = time.time() - index_start_time
            logger.info(f"Added {len(embeddings)} embeddings to FAISS index, time: {index_elapsed_time:.2f}s")

            start_idx = len(self.texts)
            for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
                text = chunk["content"]
                filename = chunk.get("file_name", "unknown")
                self.texts.append(text)
                chunk_id = hashlib.md5(text.encode("utf-8")).hexdigest()
                chunk_metadata = {
                    "file_name": filename,
                    "chunk_id": chunk_id,
                    "chunk_index": start_idx + i,
                    "metadata": chunk.get("metadata", {"tables": [], "header": None, "type": "text" if not markdown else "markdown", "sections": [], "images": []}),
                    "type": "text" if not markdown else "markdown",
                    "section": chunk.get("section", None),
                    "upload_timestamp": chunk.get("upload_timestamp", time.strftime("%Y-%m-%d %H:%M:%S"))
                }
                self.chunk_metadata.append(chunk_metadata)
                logger.debug(f"Stored chunk {i+1}/{len(chunks)} for {filename}, chunk_id: {chunk_id}")
                if filename not in self.file_data:
                    self.file_data[filename] = chunk.get("metadata", {"tables": [], "header": None, "type": "text" if not markdown else "markdown", "sections": [], "images": []})

            logger.debug(f"Saving index, texts, and metadata")
            save_start_time = time.time()
            self.save()
            save_elapsed_time = time.time() - save_start_time
            logger.info(f"Saved index, texts, and metadata, time: {save_elapsed_time:.2f}s")

            total_elapsed_time = time.time() - start_time
            logger.info(f"Successfully added {len(embeddings)} vectors for chunks, total time: {total_elapsed_time:.2f}s")
            return len(embeddings)
        except Exception as e:
            logger.error(f"Failed to add text chunks: {e}", exc_info=True)
            return 0

    def get_file_metadata(self, file_path):
        filename = os.path.basename(file_path)
        return self.file_data.get(filename, {})

    def delete_file(self, file_path):
        filename = os.path.basename(file_path)
        if filename not in self.file_data:
            logger.warning(f"File {filename} not found")
            return

        indices_to_remove = [
            i for i, meta in enumerate(self.chunk_metadata)
            if meta["file_name"] == filename
        ]

        if not indices_to_remove:
            return

        new_texts = []
        new_metadata = []
        for i, (text, meta) in enumerate(zip(self.texts, self.chunk_metadata)):
            if i not in indices_to_remove:
                new_texts.append(text)
                new_metadata.append(meta)

        self.index = faiss.IndexHNSWFlat(self.dimension, 32)
        self.index.hnsw.efConstruction = 200
        self.index.hnsw.efSearch = 100
        if new_texts:
            embeddings = self.encode_batch(new_texts)
            self.index.add(embeddings)

        self.texts = new_texts
        self.chunk_metadata = new_metadata
        del self.file_data[filename]

        self.save()
        logger.info(f"Deleted file: {filename}")

    def search(self, query, k=50, query_type="general"):
        """Search for most similar documents"""
        try:
            if self.index.ntotal == 0:
                logger.warning("Vector store is empty, cannot search")
                return []
                
            query = self.clean_text(query)
            
            try:
                lang = langdetect.detect(query)
            except:
                lang = "en"
                
            if lang in ['zh-cn', 'zh', 'ja', 'ko']:
                model_to_use = self.chinese_model
                model_name = "BAAI/bge-large-zh-v1.5"
                logger.debug(f"Using Chinese model for query: {query}")
            else:
                model_to_use = self.english_model
                model_name = "BAAI/bge-m3"
                logger.debug(f"Using English model for query: {query}")
            
            logger.info(f"Searching for query: {query} (k={k}, lang={lang}, query_type={query_type}, model={model_name})")
            
            query_embedding = model_to_use.encode([query], normalize_embeddings=query_type != "title_generation")[0]
            
            distances, indices = self.index.search(np.array([query_embedding], dtype=np.float32), k)
            
            keyword_weight = 0.4 if query_type == "technical" else 0.2 if query_type == "title_generation" else 0.3
            try:
                from jieba import analyse
                keywords = analyse.extract_tags(query, topK=10) if lang == "zh-cn" else query.split()
            except ImportError:
                logger.warning("jieba module not found, falling back to simple word splitting")
                keywords = [word for word in query.split() if len(word) > 1] if lang == "zh-cn" else query.split()
        
            contexts = []
            seen_contents = set()
            
            for i, idx in enumerate(indices[0]):
                if idx == -1 or idx >= len(self.texts):
                    continue
                meta = self.chunk_metadata[idx]
                content = self.texts[idx]
                if content in seen_contents:
                    continue
                match_count = sum(1 for keyword in keywords if keyword.lower() in content.lower())
                similarity_score = 1.0 - min(distances[0][i], 1.0)
                keyword_score = match_count / max(len(keywords), 1)
                relevance_score = (1 - keyword_weight) * similarity_score + keyword_weight * keyword_score

                table_content = ""
                if meta["metadata"].get("tables"):
                    table_content = "\n\n".join([table["text"] for table in meta["metadata"]["tables"][:1]])

                contexts.append({
                    "content": content,
                    "file_name": meta["file_name"],
                    "chunk_id": meta["chunk_id"],
                    "distance": float(distances[0][i]),
                    "relevance_score": relevance_score,
                    "metadata": meta.get("metadata", {}),
                    "type": meta.get("type", "text"),
                    "section": meta.get("section", "Unknown"),
                    "upload_timestamp": meta.get("upload_timestamp", "Unknown"),
                    "table_content": table_content
                })
                seen_contents.add(content)

            contexts = sorted(contexts, key=lambda x: x["relevance_score"], reverse=True)[:k]
            logger.info(f"Retrieved {len(contexts)} contexts for query")
            return contexts
            
        except Exception as e:
            logger.error(f"Error searching: {e}", exc_info=True)
            return []

    def search_across_files(self, query, k=5, file_names=None, query_type="general"):
        logger.info(f"Search across files: query={query}, k={k}, file_names={file_names}, query_type={query_type}")
        if self.index.ntotal == 0:
            logger.warning("Index is empty")
            return []

        try:
            try:
                lang = langdetect.detect(query)
            except Exception as e:
                logger.warning(f"Language detection failed: {str(e)}, defaulting to English")
                lang = "en"
            model = self.chinese_model if lang == "zh-cn" else self.english_model
            model_name = "BAAI/bge-large-zh-v1.5" if lang == "zh-cn" else "BAAI/bge-m3"
            logger.info(f"Using {model_name} for query embedding")

            query_embedding = model.encode([query], normalize_embeddings=query_type != "title_generation")[0]
            distances, indices = self.index.search(np.array([query_embedding]), min(k * 5, self.index.ntotal))

            contexts = []
            seen_files = set()
            keyword_weight = 0.4 if query_type == "technical" else 0.2 if query_type == "title_generation" else 0.3
            try:
                from jieba import analyse
                keywords = analyse.extract_tags(query, topK=10) if lang == "zh-cn" else query.split()
            except ImportError:
                logger.warning("jieba module not found, falling back to simple word splitting")
                keywords = [word for word in query.split() if len(word) > 1] if lang == "zh-cn" else query.split()

            for i, idx in enumerate(indices[0]):
                if idx == -1 or idx >= len(self.texts):
                    continue
                meta = self.chunk_metadata[idx]
                fname = meta["file_name"]
                if file_names and fname not in file_names:
                    continue
                if fname in seen_files:
                    continue
                content = self.texts[idx]
                match_count = sum(1 for keyword in keywords if keyword.lower() in content.lower())
                similarity_score = 1.0 - min(distances[0][i], 1.0)
                keyword_score = match_count / max(len(keywords), 1)
                relevance_score = (1 - keyword_weight) * similarity_score + keyword_weight * keyword_score

                table_content = ""
                if meta["metadata"].get("tables"):
                    table_content = "\n\n".join([table for table in meta["metadata"]["tables"][:1]])

                contexts.append({
                    "content": content,
                    "file_name": fname,
                    "chunk_id": meta["chunk_id"],
                    "distance": float(distances[0][i]),
                    "relevance_score": relevance_score,
                    "metadata": meta.get("metadata", {}),
                    "type": meta.get("type", "text"),
                    "section": meta.get("section", "Unknown"),
                    "upload_timestamp": meta.get("upload_timestamp", "Unknown"),
                    "table_content": table_content
                })
                seen_files.add(fname)
                if len(seen_files) >= k:
                    break

            contexts = sorted(contexts, key=lambda x: x["relevance_score"], reverse=True)[:k]
            logger.info(f"Retrieved {len(contexts)} contexts across files")
            return contexts

        except Exception as e:
            logger.error(f"Error searching across files: {e}", exc_info=True)
            return []
if __name__ == "__main__":
    import traceback
    sys.stdout.flush()
    sys.stderr.flush()
    print("Script starting")
    logger.info(f"enterprise_rag_optimized.py invoked, version: {SCRIPT_VERSION}")

    try:
        test_log_writability()
        logger.info("Log writability test passed")

        cache_dir = "/mnt/nfs_doc_cache"
        logger.info(f"Checking cache directory: {cache_dir}")
        print(f"Checking cache dir: {cache_dir}")
        if os.path.exists(cache_dir) and os.access(cache_dir, os.W_OK):
            logger.info(f"Cache directory {cache_dir} exists and is writable")
        else:
            logger.error(f"Cache directory {cache_dir} is not accessible or writable")
            raise FileNotFoundError(f"Cache directory {cache_dir} inaccessible")

        logger.debug("Checking dependencies")
        print("Checking dependencies")
        check_dependencies()
        logger.info("Dependencies checked successfully")

        logger.info("Checking system resources before connecting to Ray")
        print("Checking system resources")
        log_system_resources()

        # Ray cluster connection (adopted from bkp3)
        ray_host, ray_port = "192.168.0.125", 6379
        print(f"Testing network to {ray_host}:{ray_port}")
        test_network_connectivity(ray_host, ray_port)

        logger.info("Connecting to Ray cluster")
        print("Connecting to Ray cluster")
        for attempt in range(3):
            try:
                ray.init(address=f"{ray_host}:{ray_port}", namespace="vector_store", ignore_reinit_error=True)
                if ray.is_initialized():
                    nodes = ray.nodes()
                    logger.info(f"Connected to Ray cluster, nodes: {len(nodes)}")
                    print(f"Ray connected, nodes: {len(nodes)}")
                    break
                else:
                    logger.warning(f"Ray initialization failed on attempt {attempt + 1}")
            except Exception as e:
                logger.error(f"Ray connection attempt {attempt + 1} failed: {e}", exc_info=True)
                time.sleep(5)
        else:
            logger.error("Failed to connect to Ray cluster after 3 attempts")
            raise RuntimeError("Ray cluster connection failed")

        # Log cluster state for debugging
        try:
            cluster_state = ray.cluster_resources()
            logger.info(f"Cluster state: {cluster_state}")
        except Exception as e:
            logger.warning(f"Failed to get cluster state: {e}")

        # Create VectorStore actors (adopted from bkp3)
        logger.info("Creating VectorStore actors")
        print("Creating VectorStore actors")
        num_actors = 3
        actors = []
        vector_stores = []

        for i in range(num_actors):
            actor_name = f"vector_store_{i}"
            for attempt in range(3):
                try:
                    actor = ray.get_actor(actor_name, namespace="vector_store")
                    health = ray.get(actor.health_check.remote(), timeout=120)
                    if health.get("status") == "healthy":
                        logger.info(f"Found existing healthy actor {actor_name}: {health}")
                        actors.append(actor)
                        vector_stores.append(VectorStoreInfo(
                            actor=actor,
                            current_load=0,
                            max_capacity=5,
                            processing_time_avg=10.0
                        ))
                        break
                    else:
                        logger.warning(f"Actor {actor_name} unhealthy: {health}, attempting to kill")
                        ray.kill(actor)
                except ValueError:
                    logger.info(f"No existing actor found for {actor_name}, creating new")
                except Exception as e:
                    logger.warning(f"Failed to check actor {actor_name}: {e}")
                    try:
                        actor = ray.get_actor(actor_name, namespace="vector_store")
                        ray.kill(actor)
                    except ValueError:
                        pass

                try:
                    actor = VectorStore.options(
                        name=actor_name,
                        namespace="vector_store",
                        lifetime="detached",
                        num_cpus=16,  # Reduced to fit 32 CPUs
                        memory=32 * 1024 * 1024 * 1024
                    ).remote(
                        cache_dir="/mnt/nfs_doc_cache",
                        actor_id=i
                    )
                    health = ray.get(actor.health_check.remote(), timeout=180)
                    if health.get("status") == "healthy":
                        logger.info(f"Created new actor {actor_name}: {health}")
                        actors.append(actor)
                        vector_stores.append(VectorStoreInfo(
                            actor=actor,
                            current_load=0,
                            max_capacity=5,
                            processing_time_avg=10.0
                        ))
                        break
                    else:
                        logger.error(f"New actor {actor_name} unhealthy: {health}, retrying")
                        ray.kill(actor)
                except Exception as e:
                    logger.error(f"Failed to create actor {actor_name} on attempt {attempt + 1}: {e}", exc_info=True)
                    if attempt < 2:
                        logger.info(f"Retrying actor {actor_name} creation after 10 seconds")
                        time.sleep(10)
                    else:
                        logger.error(f"Failed to create actor {actor_name} after 3 attempts")
                        raise RuntimeError(f"Failed to create actor {actor_name}")

        logger.info(f"Connected to {len(actors)} VectorStore actor(s)")
        print(f"Connected to {len(actors)} actors")
        if len(actors) == 0:
            logger.error("No VectorStore actors created successfully")
            raise RuntimeError("Failed to create any VectorStore actors")

        # Verify actor-to-node mapping
        try:
            for i, store in enumerate(vector_stores):
                node_id = ray.get(store.actor.get_node_id.remote())
                logger.info(f"Actor {i} assigned to node {node_id}")
        except Exception as e:
            logger.warning(f"Failed to verify actor-to-node mapping: {e}")

        scheduler = TaskScheduler(vector_stores)
        logger.info("Task scheduler initialized")

        directory = "/mnt/nfs_docs"  # Updated to match your setup
        supported_extensions = {
            ".txt", ".md", ".csv", ".xlsx", ".xls", ".json",
            ".pdf", ".docx", ".jpg", ".jpeg", ".png", ".ppt", ".pptx", ".html"
        }
        file_paths = []
        try:
            for root, _, files in os.walk(directory):
                for fname in files:
                    if any(fname.lower().endswith(ext) for ext in supported_extensions):
                        file_path = os.path.join(root, fname)
                        file_paths.append(file_path)
            logger.info(f"Found {len(file_paths)} files to process")
        except Exception as e:
            logger.error(f"Error scanning directory {directory}: {e}", exc_info=True)
            raise RuntimeError(f"Directory scan failed: {e}")

        for file_path in file_paths:
            try:
                file_size = os.path.getsize(file_path)
                file_type = os.path.splitext(file_path)[1].lower().lstrip('.')
                if file_type == 'md':
                    file_type = 'markdown'
                priority = max(1, int(file_size / (1024 * 1024)) + 1)
                chunk_count = max(1, file_size // 1000)
                task = Task(
                    file_path=file_path,
                    chunk_count=chunk_count,
                    file_type=file_type,
                    priority=priority
                )
                scheduler.add_task(task)
                logger.debug(f"Added task: {file_path}, priority: {priority}")
            except Exception as e:
                logger.error(f"Error processing {file_path}: {e}", exc_info=True)

        # Process tasks
        futures = []
        processing_times = []
        processed_files = set()
        while not scheduler.task_queue.empty() or futures:
            try:
                assignment = scheduler.assign_task()
                if assignment:
                    actor, task = assignment
                    logger.info(f"Submitting task: {task.file_path} to actor_{actor.actor_id}")
                    future = actor.add_document.remote(task.file_path, file_type=task.file_type)
                    futures.append((future, task, actor.actor_id))
                else:
                    time.sleep(1)

                if futures:
                    ready, remaining = ray.wait(futures, num_returns=1, timeout=600)
                    futures = remaining
                    for future, task, actor_id in ready:
                        try:
                            vector_count = ray.get(future)
                            end_time = time.time()
                            processing_time = end_time - start_time
                            processing_times.append(processing_time)
                            processed_files.add(task.file_path)
                            scheduler.complete_task(actor_id)
                            logger.info(f"Completed {task.file_path}, vectors: {vector_count}, time: {processing_time:.2f}s")
                        except Exception as e:
                            logger.error(f"Error processing {task.file_path}: {e}", exc_info=True)

                for store in vector_stores:
                    health = ray.get(store.actor.health_check.remote(), timeout=600)
                    if health.get("status") != "healthy":
                        logger.warning(f"Actor_{store.actor.actor_id} unhealthy: {health.get('error', 'Unknown error')}")
            except Exception as e:
                logger.error(f"Error in main loop: {e}", exc_info=True)
                time.sleep(5)

        total_time = time.time() - start_time
        avg_processing_time = sum(processing_times) / len(processing_times) if processing_times else 0
        summary = {
            "total_files_processed": len(processed_files),
            "total_time_seconds": total_time,
            "average_processing_time_seconds": avg_processing_time,
            "total_vectors": sum(ray.get([store.actor.get_vector_count.remote()]) for store in vector_stores),
            "timestamp": time.strftime("%Y-%m-%d %H:%M:%S")
        }
        logger.info(f"Processing summary: {json.dumps(summary, indent=2)}")

        try:
            with open(os.path.join(cache_dir, "processing_summary.json"), "w", encoding="utf-8") as f:
                json.dump(summary, f, indent=2, ensure_ascii=False)
            logger.info("Processing complete, shutting down Ray")
        except Exception as e:
            logger.error(f"Failed to write processing summary: {e}", exc_info=True)

    except Exception as e:
        logger.error(f"Script execution failed: {e}", exc_info=True)
        print(f"Error: {e}")
        traceback.print_exc()
        sys.exit(1)

    finally:
        try:
            if ray.is_initialized():
                for actor in actors:
                    try:
                        logger.info(f"Requesting graceful shutdown for actor_{actor.actor_id}")
                    except Exception as e:
                        logger.warning(f"Failed to shut down actor_{actor.actor_id}: {e}")
                ray.shutdown()
                logger.info("Ray cluster connection shut down")
                print("Ray shutdown complete")
        except Exception as e:
            logger.error(f"Error during cleanup: {e}", exc_info=True)

        for handler in logger.handlers:
            handler.flush()

        logger.info(f"Script execution completed, version: {SCRIPT_VERSION}")
        print("Script execution completed")
        sys.stdout.flush()
        sys.stderr.flush()
